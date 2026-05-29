import pandas as pd
import numpy as np
from matplotlib.pyplot import plot, show, title, xlabel, ylabel, savefig
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
import os

# Load CSV files
Inflation = pd.read_csv('./CSV/US CPI.csv')
Inflation['Yearmon'] = pd.to_datetime(Inflation['Yearmon'])

OilPrices = pd.read_csv('./CSV/crude-oil-price.csv')
OilPrices['date'] = pd.to_datetime(OilPrices['date'])

Inflationcsv = './CSV/US CPI.csv'
Oilcsv = './CSV/crude-oil-price.csv'

print(OilPrices.loc[0:5, ["date", "price"]])

# Assign columns
date2 = Inflation["Yearmon"]
inflation_rate = Inflation["CPI"]

date = OilPrices["date"]
price = OilPrices["price"]
percentage_change = OilPrices["percentChange"]

# Inflation statistics
highest_inflation = inflation_rate.max()
lowest_inflation = inflation_rate.min()

highest_inflation_date = date2[inflation_rate.idxmax()]
lowest_inflation_date = date2[inflation_rate.idxmin()]

# Oil price statistics
highest_price = price.max()
lowest_price = price.min()

highest_date = date[price.idxmax()]
lowest_date = date[price.idxmin()]

print(f"Highest price: ${highest_price} on {highest_date}")
print(f"Lowest price: ${lowest_price} on {lowest_date}")


# Oil Price Graph
plt.figure()

plot(date, price)

title("Oil Price Over Time", fontsize=20)
xlabel('Date', fontsize=14)
ylabel('Price', fontsize=14)

savefig("oil_price_graph.png", dpi=100, bbox_inches='tight')


# Percentage Change Graph
plt.figure()

plot(date, percentage_change)

title("Daily Percentage Change in Oil Price", fontsize=20)
xlabel('Date', fontsize=14)
ylabel('Percent Change', fontsize=14)

savefig("oil_price_percentage_change.png", dpi=100, bbox_inches='tight')


# Inflation Graph
plt.figure()

plot(date2, inflation_rate)

title("Inflation Rate Over Time", fontsize=20)
xlabel('Date', fontsize=14)
ylabel('Inflation Rate', fontsize=14)

savefig("inflation_rate_graph.png", dpi=100, bbox_inches='tight')


# Filter Inflation Data

# Remove timezone information if it exists
Inflation["Yearmon"] = pd.to_datetime(
    Inflation["Yearmon"]
).dt.tz_localize(None)

OilPrices["date"] = pd.to_datetime(
    OilPrices["date"]
).dt.tz_localize(None)

# Get earliest oil price date
start_date = OilPrices["date"].min()

# Filter CPI data to start from same date
filtered_inflation = Inflation[
    Inflation["Yearmon"] >= start_date
]

filtered_dates = filtered_inflation["Yearmon"]
filtered_cpi = filtered_inflation["CPI"]


# Combined Graph
plt.figure(figsize=(10, 6))

plot(date, price, label="Oil Price")
plot(filtered_dates, filtered_cpi, label="Inflation Rate")

title("Inflation Rate and Oil Price Over Time", fontsize=20)
xlabel("Date", fontsize=14)
ylabel("Values", fontsize=14)

plt.legend()

savefig("combined_graph.png", dpi=100, bbox_inches='tight')


# Create PowerPoint
pr1 = Presentation()

slide = pr1.slides.add_slide(pr1.slide_layouts[5])
slide2 = pr1.slides.add_slide(pr1.slide_layouts[5])
slide3 = pr1.slides.add_slide(pr1.slide_layouts[5])
slide4 = pr1.slides.add_slide(pr1.slide_layouts[5])
slide5 = pr1.slides.add_slide(pr1.slide_layouts[5])
slide6 = pr1.slides.add_slide(pr1.slide_layouts[5])
slide7 = pr1.slides.add_slide(pr1.slide_layouts[5])

# Slide 1
title = slide.shapes.title
title.text = "Oil Price Over Time"


# Slide 2 - Oil Price Graph
title2 = slide2.shapes.title
title2.text = "Highest and Lowest Oil Price"

left = top = 0

pic = slide2.shapes.add_picture(
    "oil_price_graph.png",
    left,
    top,
    width=9144000,
    height=6858000
)


# Slide 3 - Oil Statistics
title3 = slide3.shapes.title
title3.text = "Oil Price Statistics"

left = 1000000
top = 2500000
width = 7000000
height = 3000000

textbox = slide3.shapes.add_textbox(left, top, width, height)

text_frame = textbox.text_frame
text_frame.word_wrap = True

p1 = text_frame.paragraphs[0]
p1.text = (
    f"Highest price: ${highest_price} "
    f"on {highest_date.strftime('%Y-%m-%d')}"
)
p1.font.size = 304800

p2 = text_frame.add_paragraph()
p2.text = (
    f"\nLowest price: ${lowest_price} "
    f"on {lowest_date.strftime('%Y-%m-%d')}"
)
p2.font.size = 304800


# Slide 4 - Percentage Change
title4 = slide4.shapes.title
title4.text = "Oil Price Percentage Change"

left = top = 0

pic2 = slide4.shapes.add_picture(
    "oil_price_percentage_change.png",
    left,
    top,
    width=9144000,
    height=6858000
)


# Slide 5 - Inflation Graph
title5 = slide5.shapes.title
title5.text = "US CPI Over Time"

left = top = 0

pic3 = slide5.shapes.add_picture(
    "inflation_rate_graph.png",
    left,
    top,
    width=9144000,
    height=6858000
)


# Slide 6 - CPI Statistics
title6 = slide6.shapes.title
title6.text = "Highest and Lowest Inflation CPI"

left = 1000000
top = 2500000
width = 7000000
height = 3000000

textbox2 = slide6.shapes.add_textbox(left, top, width, height)

text_frame2 = textbox2.text_frame
text_frame2.word_wrap = True

p3 = text_frame2.paragraphs[0]
p3.text = (
    f"Highest CPI: {highest_inflation} "
    f"on {highest_inflation_date.strftime('%Y-%m-%d')}"
)
p3.font.size = 304800

p4 = text_frame2.add_paragraph()
p4.text = (
    f"\nLowest CPI: {lowest_inflation} "
    f"on {lowest_inflation_date.strftime('%Y-%m-%d')}"
)
p4.font.size = 304800

# Slide 7 - Combined Graph
title7 = slide7.shapes.title
title7.text = "Inflation Rate vs Oil Price"

left = top = 0

pic4 = slide7.shapes.add_picture(
    "combined_graph.png",
    left,
    top,
    width=9144000,
    height=6858000
)


# Save PowerPoint
pr1.save("Bogdans_Presentation.pptx")

# Open PowerPoint (optional - comment out if you don't want it to auto-open)
try:
    os.startfile("Bogdans_Presentation.pptx")
except Exception as e:
    print(f"Could not open PowerPoint: {e}")

print("klar!")